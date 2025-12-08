"""
Create a PowerPoint presentation summarizing error analyses per variable.

Each slide contains:
- Top-left: variable_type, variable_name and stats (VRMSE, Std error, Max error sample, Max error timepoint)
- Bottom-left: three histograms side-by-side (initial timepoint, up to N timepoints, all timepoints)
- Right: error over time, worst sample, #2 worst, #3 worst, median sample (true vs pred)

This script reuses utilities from plot_error_histogram, plot_error_over_time, and plot_dataset.
"""

from pathlib import Path
from typing import Dict, List, Tuple
import argparse
import logging
import sys

import h5py
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

import filepaths

from common import save_figure
from plot_dataset import (
    calculate_scaling_dict,
    calculate_normalization_stats,
    create_multi_axis_figure,
)
from plot_error_histogram import (
    calculate_variable_errors,
    get_variable_names,
    collect_error_data,
    _draw_error_histogram,
    make_pandas_dataframe,
)
from plot_error_over_time import _draw_error_time_series


def _sanitize(name: str) -> str:
    return name.replace('/', '_').replace(':', '_').replace(' ', '_')


def _errors_for_time_selection(err: np.ndarray, n_tp: int) -> np.ndarray:
    # Accepts (n_samples, T) or (n_samples,)
    e = np.asarray(err)
    if e.ndim == 1:
        return e.reshape(-1)
    if n_tp <= 0 or n_tp >= e.shape[1]:
        return e.reshape(-1)
    return e[:, :n_tp].reshape(-1)


def _compute_worst_and_median_indices(
    errors: np.ndarray,
    top_k: int = 3,
) -> Tuple[List[int], int, np.ndarray]:
    """Compute worst and median sample indices based on per-sample VRMSE.

    Args:
        errors: Array shaped (n_samples, T) for time series or (n_samples,) for parameters.
        top_k: Number of worst samples to return (default: 3).

    Returns:
        (worst_indices, median_index, vrmse_per_sample)
    """
    e = np.asarray(errors)
    if e.ndim == 1:
        # Single timepoint per sample; use absolute error as VRMSE
        vrmse = np.abs(e)
    else:
        # Compute per-sample VRMSE across timepoints (last dim)
        # Ensure shape (n_samples, T)
        if e.ndim > 2:
            # Flatten trailing dimensions into time
            n = e.shape[0]
            e = e.reshape(n, -1)
        vrmse = np.sqrt(np.nanmean(e**2, axis=1))

    # Handle NaNs by treating them as very small (to not be selected as worst)
    vrmse_clean = np.nan_to_num(vrmse, nan=-np.inf)

    n = vrmse_clean.shape[0]
    k = min(top_k, n)
    order = np.argsort(vrmse_clean)
    worst_indices = order[-k:][::-1].tolist()
    median_index = order[n // 2].item()
    return worst_indices, median_index, vrmse


def _save_hist_figure(
    errors_1d: np.ndarray,
    bins: int,
    error_range: float,
    title: str,
    out_path: Path,
    dpi: int,
    *,
    figsize: Tuple[float, float],
) -> Path:
    fig, ax = plt.subplots(figsize=figsize)
    _draw_error_histogram(
        ax,
        errors_1d,
        bins=bins,
        error_range=error_range,
        title=title,
        show_legend=False,
        annotate_text=True,
        title_fontsize=9,
        label_fontsize=8,
    )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    save_figure(fig, out_path.parent, out_path.stem, 'png', dpi=dpi)
    plt.close(fig)
    return out_path


def _save_error_over_time_figure(
    err: np.ndarray,
    title: str,
    out_path: Path,
    max_y: float,
    dpi: int,
    *,
    figsize: Tuple[float, float],
) -> Path:
    fig, ax = plt.subplots(figsize=figsize)
    _draw_error_time_series(ax, err, title=title, max_y=max_y)
    save_figure(fig, out_path.parent, out_path.stem, 'png', dpi=dpi)
    plt.close(fig)
    return out_path


def _save_prediction_figure(
    hdf5_file: h5py.File,
    scaling_dict: Dict[str, Dict[str, Tuple[float, float]]],
    normalization_stats: Dict[str, Dict[str, Tuple[float, float]]],
    context_key: str,
    sample_idx: int,
    var_type: str,
    var_name: str,
    description: str,
    out_path: Path,
    dpi: int,
    *,
    figsize: Tuple[float, float],
) -> Path:
    plot_specs = [{
        'variable_type': var_type,
        'variable_names': [var_name],
    }]
    fig, _ = create_multi_axis_figure(
        hdf5_file,
        scaling_dict,
        context_key,
        sample_idx,
        plot_specs,
        description_text=description,
        figsize=figsize,
        normalization_stats=normalization_stats,
    )
    save_figure(fig, out_path.parent, out_path.stem, 'png', dpi=dpi)
    plt.close(fig)
    return out_path


def generate_ppt(
    hdf5_file: h5py.File,
    context_key: str,
    output_pptx: Path,
    *,
    bins: int = 50,
    error_range: float = 0.6,
    mid_timepoints: int = 10,
    max_y: float = 0.8,
    dpi: int = 200,
    test_mode: bool = False,
) -> None:
    output_pptx = Path(output_pptx)
    assets_dir = output_pptx.parent / (output_pptx.stem + "_assets")
    assets_dir.mkdir(parents=True, exist_ok=True)

    # Stats and errors
    normalization_stats = calculate_normalization_stats(hdf5_file, 'train')
    scaling_dict = calculate_scaling_dict(hdf5_file, 'train')
    variable_errors = calculate_variable_errors(hdf5_file, context_key, normalization_stats)
    variable_names: Dict[str, List[str]] = {}
    variable_names = get_variable_names(hdf5_file, normalization_stats, variable_names)
    all_variable_names, error_vectors = collect_error_data(variable_errors, variable_names)
    df_stats = make_pandas_dataframe(variable_names, variable_errors)

    # Optionally restrict for test mode
    if test_mode:
        all_variable_names = all_variable_names[: min(4, len(all_variable_names))]
        error_vectors = error_vectors[: len(all_variable_names)]
        logging.info(f"Test mode: limiting to {len(all_variable_names)} variables")

    prs = Presentation("utils/plots/template.pptx") # TODO: add an empty template file and a CLI option to use a custom template
    prs.core_properties.title = f"Error Analysis - {context_key}"

    # Slide geometry in inches
    slide_w_in = prs.slide_width / Inches(1)
    slide_h_in = prs.slide_height / Inches(1)

    # Layout params (inches)
    margin_l = 0.3
    margin_r = 0.3
    margin_t = 0.2
    margin_b = 0.3
    col_gap = 0.25
    block_gap = 0.12

    # Column widths
    right_col_frac = 0.4  # allocate 40% of width to right column
    right_col_w = max(3.2, min(slide_w_in * right_col_frac, slide_w_in - margin_l - margin_r - 3.0))
    left_col_w = slide_w_in - margin_l - margin_r - col_gap - right_col_w

    # Text box height
    text_h = max(0.9, min(1.3, slide_h_in * 0.16))

    # Histogram target aspect (match our default 3.2x2.6)
    hist_aspect = 3.2 / 2.6  # width / height

    for name, err in zip(all_variable_names, error_vectors):
        var_type, var_name = [s.strip() for s in name.split(':', 1)]
        logging.info(f"Adding slide for {var_type}/{var_name}")

        # Find stats row
        row = df_stats[(df_stats['Variable type'] == var_type) & (df_stats['Variable'] == var_name)]
        stats_text = "Stats not found"
        if not row.empty:
            r = row.iloc[0]
            stats_text = (
                f"Type: {var_type}\n"
                f"Name: {var_name}\n"
                f"VRMSE: {r['VRMSE']:.4f}\n"
                f"Std error: {r['Std error']:.4f}\n"
                f"Max error sample: {r['Max error sample']:.4f}\n"
                f"Max error timepoint: {r['Max error timepoint']:.4f}"
            )

        # Prepare images
        base = _sanitize(f"{var_type}_{var_name}")
        # Histograms data
        e_init = _errors_for_time_selection(err, 1)
        e_mid = _errors_for_time_selection(err, mid_timepoints)
        e_all = _errors_for_time_selection(err, -1)

        # Compute histogram sizes to fit left column
        # First compute width-driven size
        w_hist = (left_col_w - 2 * block_gap) / 3.0
        h_hist = w_hist / hist_aspect
        # Cap by available vertical space
        max_hist_h = slide_h_in - margin_t - text_h - margin_b - block_gap
        if h_hist > max_hist_h:
            h_hist = max_hist_h
            w_hist = h_hist * hist_aspect
            # Ensure three histograms still fit horizontally; if not, shrink to fit
            total_w = 3 * w_hist + 2 * block_gap
            if total_w > left_col_w:
                w_hist = (left_col_w - 2 * block_gap) / 3.0
                h_hist = w_hist / hist_aspect

        # Error-over-time and predictions sizes (stacked equally in right column)
        # Worst list (up to 3) + median + error-over-time
        worst_idxs, median_idx, vrmse_per_sample = _compute_worst_and_median_indices(err)
        n_right_tiles = 1 + len(worst_idxs) + 1  # eot + worsts + median
        w_right = right_col_w
        # Compute equal heights with gaps
        h_right = (slide_h_in - margin_t - margin_b - (n_right_tiles - 1) * block_gap) / n_right_tiles
        # Reasonable cap
        h_right = max(0.9, min(1.8, h_right))

        # File paths for images
        p_init = assets_dir / f"{base}_hist_init.png"
        p_mid = assets_dir / f"{base}_hist_up_to_{mid_timepoints}.png"
        p_all = assets_dir / f"{base}_hist_all.png"
        p_eot = assets_dir / f"{base}_error_over_time.png"

        # Save figures with figsize matching insertion sizes to avoid distortion
        _save_hist_figure(e_init, bins, error_range, "Initial TP", p_init, dpi, figsize=(2*w_hist, 2*h_hist))
        _save_hist_figure(e_mid, bins, error_range, f"Up to {mid_timepoints}", p_mid, dpi, figsize=(2*w_hist, 2*h_hist))
        _save_hist_figure(e_all, bins, error_range, "All TPs", p_all, dpi, figsize=(2*w_hist, 2*h_hist))

        _save_error_over_time_figure(err, "|error| over time", p_eot, max_y=max_y, dpi=dpi, figsize=(2*w_right, 2*h_right))

        # Worst and median samples
        pics_pred: List[Path] = []
        # Up to 3 worst
        for rank, idx in enumerate(worst_idxs, start=1):
            desc = f"Worst #{rank} (VRMSE: {vrmse_per_sample[idx]:.4f})"
            p = assets_dir / f"{base}_worst_{rank}.png"
            _save_prediction_figure(
                hdf5_file, scaling_dict, normalization_stats, context_key,
                idx, var_type, var_name, desc, p, dpi, figsize=(2*w_right, 2*h_right)
            )
            pics_pred.append(p)
        # Median
        p_med = assets_dir / f"{base}_median.png"
        _save_prediction_figure(
            hdf5_file, scaling_dict, normalization_stats, context_key,
            median_idx, var_type, var_name, f"Median (VRMSE: {vrmse_per_sample[median_idx]:.4f})",
            p_med, dpi, figsize=(2*w_right, 2*h_right)
        )
        pics_pred.append(p_med)

        # Compose slide
        slide = prs.slides.add_slide(prs.slide_layouts[10])  # blank

        # Text box (top-left)
        left = Inches(margin_l); top = Inches(margin_t); width = Inches(left_col_w); height = Inches(text_h)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stats_text
        p.font.size = Pt(14)

        # Bottom-left histograms (side-by-side)
        y_hist_top = margin_t + text_h + block_gap + 0.8
        slide.shapes.add_picture(str(p_init), Inches(margin_l + 0 * (w_hist + block_gap)), Inches(y_hist_top), width=Inches(w_hist), height=Inches(h_hist))
        slide.shapes.add_picture(str(p_mid), Inches(margin_l + 1 * (w_hist + block_gap)), Inches(y_hist_top), width=Inches(w_hist), height=Inches(h_hist))
        slide.shapes.add_picture(str(p_all), Inches(margin_l + 2 * (w_hist + block_gap)), Inches(y_hist_top), width=Inches(w_hist), height=Inches(h_hist))

        # Right column stacked images
        x_right = slide_w_in - margin_r - right_col_w
        y = margin_t
        slide.shapes.add_picture(str(p_eot), Inches(x_right), Inches(y), width=Inches(w_right), height=Inches(h_right))
        y += h_right + block_gap
        for pic in pics_pred:
            slide.shapes.add_picture(str(pic), Inches(x_right), Inches(y), width=Inches(w_right), height=Inches(h_right))
            y += h_right + block_gap

    prs.save(str(output_pptx))
    logging.info(f"Saved presentation to {output_pptx}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Create a PowerPoint of error analyses per variable")
    parser.add_argument('--dataset_path', type=str, required=True, help='Path to HDF5 dataset (local or ML artifacts)')
    parser.add_argument('--context', type=str, default='common_test', help='Context to analyze (default: common_test)')
    parser.add_argument('--output_pptx', type=str, required=True, help='Output .pptx path')
    parser.add_argument('--bins', type=int, default=50, help='Histogram bins (default: 50)')
    parser.add_argument('--error-range', type=float, default=0.6, dest='error_range', help='Histogram x-range half-width (default: 0.6)')
    parser.add_argument('--mid-timepoints', type=int, default=10, dest='mid_timepoints', help='Timepoints used in the middle histogram (default: 10)')
    parser.add_argument('--max-y', type=float, default=0.8, dest='max_y', help='Max y for error-over-time plots (default: 0.8)')
    parser.add_argument('--dpi', type=int, default=200, help='DPI for saved figures (default: 200)')
    parser.add_argument('--test', action='store_true', dest='test_mode', help='Limit to a few variables for a quick run')

    if '--help' in sys.argv or '-h' in sys.argv:
        print("Usage example:")
        print("create_ppt.py --dataset_path mlflow-artifacts:/689985610175568372/5b7ce6ec47694a5f9661e49ee1be98d0/artifacts/dataset.hdf5 --context common_test --output_pptx output/error_analysis.pptx --test")

    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

    try:
        with h5py.File(filepaths.filepath_from_local_or_ml_artifacts(args.dataset_path), 'r') as f:
            generate_ppt(
                f,
                args.context,
                Path(args.output_pptx),
                bins=args.bins,
                error_range=args.error_range,
                mid_timepoints=args.mid_timepoints,
                max_y=args.max_y,
                dpi=args.dpi,
                test_mode=args.test_mode,
            )
        print(f"Presentation saved to {args.output_pptx}")
    except Exception as e:
        logging.error(f"Error creating presentation: {e}")
        raise